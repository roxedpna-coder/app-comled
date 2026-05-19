from pptx import Presentation
from pptx.util import Pt
import json
import os
import ast
from PIL import Image

# -------------------------
# CARGAR JSON
# -------------------------

with open("entradas/datos.json", "r", encoding="utf-8") as archivo:
    datos = json.load(archivo)

# -------------------------
# ABRIR PLANTILLA
# -------------------------

ppt = Presentation("plantillas/FICHA_COMLED_MASTER.pptx")

# -------------------------
# LIMPIAR VALORES
# -------------------------

def limpiar_valor(valor):
    if isinstance(valor, list):
        return " / ".join(str(v).capitalize() for v in valor)

    valor = str(valor).strip()

    if valor.startswith("[") and valor.endswith("]"):
        try:
            lista = ast.literal_eval(valor)
            if isinstance(lista, list):
                return " / ".join(str(v).capitalize() for v in lista)
        except:
            pass

    return valor

# -------------------------
# ICONO DE INSTALACIÓN
# -------------------------

def obtener_icono_instalacion(datos):
    instalacion = str(datos.get("instalacion", "")).lower().strip()

    # Prioridad absoluta: campo instalación
    if "empotr" in instalacion or "downlight" in instalacion:
        return "imagenes/downlights.png"

    if "carril" in instalacion:
        return "imagenes/de-carril.png"

    if "superficie" in instalacion:
        return "imagenes/de-superficie.png"

    if "pared" in instalacion or "muro" in instalacion:
        return "imagenes/de-pared.png"

    if "sobremesa" in instalacion or "mesa" in instalacion:
        return "imagenes/de-sobremesa.png"

    if "lineal" in instalacion or "perfil" in instalacion:
        return "imagenes/lineales.png"

    if "colg" in instalacion or "suspend" in instalacion:
        return "imagenes/colgantes.png"

    if "proyector" in instalacion:
        return "imagenes/proyectores.png"

    if "baliza" in instalacion:
        return "imagenes/balizas.png"

    if "emergencia" in instalacion:
        return "imagenes/emergencia.png"

    if (
        "señal" in instalacion
        or "senaletica" in instalacion
        or "señaletica" in instalacion
    ):
        return "imagenes/senaletica.png"

    if "integracion" in instalacion or "integración" in instalacion:
        return "imagenes/integracion.png"

    if "uplight" in instalacion:
        return "imagenes/uplights.png"

    return "imagenes/tecnicas.png"

# -------------------------
# ÓPTICA REPRESENTADA
# -------------------------

optica_grafico = datos.get("optica_grafico", "-")

# -------------------------
# REEMPLAZOS
# -------------------------

import re

def resumir_ip(valor):
    valor = str(valor).upper().strip()
    if valor.count("IP") > 1:
        partes = [p.strip() for p in valor.replace("&", "/").replace("-", "/").split("/")]
        limpias = []
        for p in partes:
            num = re.sub(r"[^\d]", "", p)
            if num: limpias.append(num)
        if limpias:
            return "IP" + "/".join(limpias)
    return valor

def resumir_rango(valor, sufijo=""):
    valor = str(valor).strip()
    numeros = re.findall(r"[\d]+", valor.replace(".", "").replace(",", ""))
    if len(numeros) > 1:
        nums = [int(n) for n in numeros if int(n) > 0]
        if nums:
            return f"{min(nums)}-{max(nums)}{sufijo}"
    return valor

# Soporte de compatibilidad para plantillas antiguas/nuevas y forzado de resumen gráfico
val_cct = str(datos.get("cct_resumido", ""))
if not val_cct or len(val_cct) > 15 or "/" in val_cct:
    datos["cct_resumido"] = resumir_rango(datos.get("cct", "-"), "K")

val_ip = str(datos.get("ip_resumido", ""))
if not val_ip or len(val_ip) > 9 or "IP" in val_ip[3:]:
    datos["ip_resumido"] = resumir_ip(datos.get("ip", "-"))

val_pot = str(datos.get("potencia_resumido", ""))
if not val_pot or len(val_pot) > 10 or "/" in val_pot:
    datos["potencia_resumido"] = resumir_rango(datos.get("potencia", "-"), "W")

val_flu = str(datos.get("flujo_resumido", ""))
if not val_flu or len(val_flu) > 12 or "/" in val_flu:
    datos["flujo_resumido"] = resumir_rango(datos.get("flujo_luminoso", "-"), "lm")

val_efi = str(datos.get("eficacia_resumido", ""))
if not val_efi or len(val_efi) > 12 or "/" in val_efi:
    datos["eficacia_resumido"] = resumir_rango(datos.get("eficacia_luminosa", "-"), "lm/W")

reemplazos = {}

for clave, valor in datos.items():
    if not clave.startswith("img_"):
        reemplazos["{{" + clave + "}}"] = limpiar_valor(valor)

reemplazos["{{optica_grafico}}"] = optica_grafico

# -------------------------
# IMÁGENES
# -------------------------

imagenes = {
    "{{img_producto}}": datos.get(
        "img_producto",
        "imagenes/producto_final.png"
    ),

    "{{img_dimensiones}}": datos.get(
        "img_dimensiones",
        "imagenes/dimensiones.png"
    ),

    "{{img_fotometria}}": datos.get(
        "img_fotometria",
        "imagenes/fotometria_generada.png"
    ),

    "{{img_instalacion}}": obtener_icono_instalacion(datos)
}

# -------------------------
# ESCALAS
# -------------------------

escalas = {
    "{{img_producto}}": 1.00,
    "{{img_dimensiones}}": 1.00,
    "{{img_fotometria}}": 1.00,
    "{{img_instalacion}}": 1.35
}

# -------------------------
# LEER TEXTO
# -------------------------

def leer_texto(shape):
    texto = ""

    try:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    texto += run.text
    except:
        pass

    return texto

# -------------------------
# FORMATEAR TEXTO
# -------------------------

def formatear_shape(shape, size=8, bold=False):
    try:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = "Arial"
                run.font.size = Pt(size)
                run.font.bold = bold
                run.font._element.set("b", "1" if bold else "0")
    except:
        pass

# -------------------------
# REEMPLAZAR TEXTO
# -------------------------

def reemplazar_texto_sin_romper_formato(shape, placeholder, valor):
    try:
        if not shape.has_text_frame:
            return

        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if placeholder in run.text:
                    run.text = run.text.replace(
                        placeholder,
                        str(valor)
                    )

        if placeholder == "{{descripcion_luminaria}}":
            formatear_shape(shape, size=8, bold=False)

        if placeholder == "{{optica_grafico}}":
            formatear_shape(shape, size=8, bold=True)

    except:
        pass

# -------------------------
# INSERTAR IMAGEN
# -------------------------

def insertar_imagen_sin_deformar(slide, shape, ruta_imagen, escala=1.0):
    left = shape.left
    top = shape.top

    box_width = shape.width
    box_height = shape.height

    with Image.open(ruta_imagen) as img:
        img_width, img_height = img.size

    ratio_img = img_width / img_height
    ratio_box = box_width / box_height

    if ratio_img > ratio_box:
        new_width = box_width
        new_height = int(box_width / ratio_img)
    else:
        new_height = box_height
        new_width = int(box_height * ratio_img)

    new_width = int(new_width * escala)
    new_height = int(new_height * escala)

    if new_width > box_width:
        new_width = box_width
        new_height = int(new_width / ratio_img)

    if new_height > box_height:
        new_height = box_height
        new_width = int(new_height * ratio_img)

    new_left = left + int((box_width - new_width) / 2)
    new_top = top + int((box_height - new_height) / 2)

    slide.shapes._spTree.remove(shape._element)

    slide.shapes.add_picture(
        ruta_imagen,
        new_left,
        new_top,
        width=new_width,
        height=new_height
    )

# -------------------------
# RECORRER SHAPES
# -------------------------

for slide in ppt.slides:
    for shape in list(slide.shapes):
        texto_shape = leer_texto(shape)

        imagen_reemplazada = False

        for placeholder, ruta_imagen in imagenes.items():
            if placeholder in texto_shape:
                if os.path.exists(ruta_imagen):
                    insertar_imagen_sin_deformar(
                        slide,
                        shape,
                        ruta_imagen,
                        escalas.get(placeholder, 1.0)
                    )
                else:
                    print("No existe la imagen:", ruta_imagen)

                imagen_reemplazada = True
                break

        if imagen_reemplazada:
            continue

        for placeholder, valor in reemplazos.items():
            if placeholder in texto_shape:
                reemplazar_texto_sin_romper_formato(
                    shape,
                    placeholder,
                    valor
                )

# -------------------------
# GUARDAR
# -------------------------

ppt.save("salidas/ficha_final.pptx")

print("Ficha generada correctamente")
